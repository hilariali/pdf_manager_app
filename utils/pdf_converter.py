import io
from pathlib import Path
import fitz  # PyMuPDF
from PIL import Image
import pandas as pd
from docx import Document
from pptx import Presentation
import zipfile
import tempfile
import os

class PDFConverter:
    def __init__(self):
        self.supported_formats = {
            'pdf_to_word': self.pdf_to_word,
            'pdf_to_excel': self.pdf_to_excel,
            'pdf_to_powerpoint': self.pdf_to_powerpoint,
            'pdf_to_images': self.pdf_to_images,
            'word_to_pdf': self.word_to_pdf,
            'excel_to_pdf': self.excel_to_pdf,
            'powerpoint_to_pdf': self.powerpoint_to_pdf,
            'images_to_pdf': self.images_to_pdf
        }
    
    def convert_file(self, uploaded_file, conversion_type):
        """Main conversion dispatcher"""
        conversion_key = conversion_type.lower().replace(' ', '_')
        
        if conversion_key in self.supported_formats:
            return self.supported_formats[conversion_key](uploaded_file)
        else:
            raise ValueError(f"Unsupported conversion type: {conversion_type}")
    
    def pdf_to_word(self, pdf_file):
        """Convert PDF to Word document"""
        doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
        word_doc = Document()
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            word_doc.add_paragraph(text)
            if page_num < len(doc) - 1:  # Don't add page break after last page
                word_doc.add_page_break()
        
        output = io.BytesIO()
        word_doc.save(output)
        output.seek(0)
        doc.close()
        
        return {
            'data': output.getvalue(),
            'filename': f"{Path(pdf_file.name).stem}.docx",
            'mime_type': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        }
    
    def pdf_to_excel(self, pdf_file):
        """Convert PDF to Excel document"""
        doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
        
        # Create a list to store all text data
        all_data = []
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            
            # Split text into lines and create rows
            lines = text.split('\n')
            for line in lines:
                if line.strip():  # Skip empty lines
                    # Simple approach: split by spaces or tabs
                    row_data = line.split()
                    all_data.append(row_data)
        
        # Create DataFrame
        if all_data:
            # Find the maximum number of columns
            max_cols = max(len(row) for row in all_data) if all_data else 1
            
            # Pad rows to have the same number of columns
            padded_data = []
            for row in all_data:
                padded_row = row + [''] * (max_cols - len(row))
                padded_data.append(padded_row)
            
            df = pd.DataFrame(padded_data)
        else:
            df = pd.DataFrame([['No data extracted']])
        
        output = io.BytesIO()
        df.to_excel(output, index=False, header=False)
        output.seek(0)
        doc.close()
        
        return {
            'data': output.getvalue(),
            'filename': f"{Path(pdf_file.name).stem}.xlsx",
            'mime_type': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        }
    
    def pdf_to_powerpoint(self, pdf_file):
        """Convert PDF to PowerPoint presentation"""
        doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
        ppt = Presentation()
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            
            # Create a new slide
            slide_layout = ppt.slide_layouts[1]  # Title and Content layout
            slide = ppt.slides.add_slide(slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = f"Page {page_num + 1}"
            
            # Add content
            content = slide.placeholders[1]
            content.text = text[:1000]  # Limit text length
        
        output = io.BytesIO()
        ppt.save(output)
        output.seek(0)
        doc.close()
        
        return {
            'data': output.getvalue(),
            'filename': f"{Path(pdf_file.name).stem}.pptx",
            'mime_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        }
    
    def pdf_to_images(self, pdf_file):
        """Convert PDF pages to images"""
        doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
        images = []
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom
            img_data = pix.tobytes("png")
            images.append(img_data)
        
        # Create ZIP file with all images
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w') as zip_file:
            for i, img_data in enumerate(images):
                zip_file.writestr(f"page_{i+1}.png", img_data)
        
        zip_buffer.seek(0)
        doc.close()
        
        return {
            'data': zip_buffer.getvalue(),
            'filename': f"{Path(pdf_file.name).stem}_images.zip",
            'mime_type': 'application/zip'
        }
    
    def word_to_pdf(self, word_file):
        """Convert Word document to PDF"""
        try:
            # Read the Word document
            doc = Document(word_file)
            
            # Create a new PDF document
            pdf_doc = fitz.open()
            
            # Extract text from Word document
            full_text = []
            for paragraph in doc.paragraphs:
                full_text.append(paragraph.text)
            
            text_content = '\n'.join(full_text)
            
            # Create PDF page with text
            page = pdf_doc.new_page()
            
            # Insert text into PDF
            text_rect = fitz.Rect(50, 50, page.rect.width - 50, page.rect.height - 50)
            page.insert_textbox(text_rect, text_content, fontsize=12, fontname="helv")
            
            output = io.BytesIO()
            pdf_doc.save(output)
            output.seek(0)
            pdf_doc.close()
            
            return {
                'data': output.getvalue(),
                'filename': f"{Path(word_file.name).stem}.pdf",
                'mime_type': 'application/pdf'
            }
            
        except Exception as e:
            raise ValueError(f"Failed to convert Word to PDF: {str(e)}")
    
    def excel_to_pdf(self, excel_file):
        """Convert Excel document to PDF"""
        try:
            # Read the Excel file
            df = pd.read_excel(excel_file)
            
            # Create a new PDF document
            pdf_doc = fitz.open()
            page = pdf_doc.new_page()
            
            # Convert DataFrame to string
            text_content = df.to_string(index=False)
            
            # Insert text into PDF
            text_rect = fitz.Rect(50, 50, page.rect.width - 50, page.rect.height - 50)
            page.insert_textbox(text_rect, text_content, fontsize=10, fontname="helv")
            
            output = io.BytesIO()
            pdf_doc.save(output)
            output.seek(0)
            pdf_doc.close()
            
            return {
                'data': output.getvalue(),
                'filename': f"{Path(excel_file.name).stem}.pdf",
                'mime_type': 'application/pdf'
            }
            
        except Exception as e:
            raise ValueError(f"Failed to convert Excel to PDF: {str(e)}")
    
    def powerpoint_to_pdf(self, ppt_file):
        """Convert PowerPoint presentation to PDF"""
        try:
            # Read the PowerPoint file
            ppt = Presentation(ppt_file)
            
            # Create a new PDF document
            pdf_doc = fitz.open()
            
            for slide_num, slide in enumerate(ppt.slides):
                # Create a new page for each slide
                page = pdf_doc.new_page()
                
                # Extract text from slide
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        slide_text.append(shape.text)
                
                text_content = '\n'.join(slide_text)
                
                # Add slide title
                title = f"Slide {slide_num + 1}\n{'='*20}\n"
                full_text = title + text_content
                
                # Insert text into PDF
                text_rect = fitz.Rect(50, 50, page.rect.width - 50, page.rect.height - 50)
                page.insert_textbox(text_rect, full_text, fontsize=12, fontname="helv")
            
            output = io.BytesIO()
            pdf_doc.save(output)
            output.seek(0)
            pdf_doc.close()
            
            return {
                'data': output.getvalue(),
                'filename': f"{Path(ppt_file.name).stem}.pdf",
                'mime_type': 'application/pdf'
            }
            
        except Exception as e:
            raise ValueError(f"Failed to convert PowerPoint to PDF: {str(e)}")
    
    def images_to_pdf(self, image_files):
        """Convert multiple images to a single PDF"""
        if not isinstance(image_files, list):
            image_files = [image_files]
        
        pdf_document = fitz.open()
        
        for image_file in image_files:
            try:
                img = Image.open(image_file)
                
                # Convert image to RGB if necessary
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                img_bytes = io.BytesIO()
                img.save(img_bytes, format='PNG')
                img_bytes.seek(0)
                
                # Create new page with image dimensions
                page = pdf_document.new_page(width=img.width, height=img.height)
                page.insert_image(page.rect, stream=img_bytes.getvalue())
                
            except Exception as e:
                print(f"Error processing image {image_file.name}: {str(e)}")
                continue
        
        output = io.BytesIO()
        pdf_document.save(output)
        output.seek(0)
        pdf_document.close()
        
        return {
            'data': output.getvalue(),
            'filename': "converted_images.pdf",
            'mime_type': 'application/pdf'
        }
